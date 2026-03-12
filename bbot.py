# Import necessary libraries
import os
import logging
import asyncio
from aiogram import Bot, Dispatcher, types
from aiogram.contrib.middlewares.logging import LoggingMiddleware
from aiogram.dispatcher import FSMContext
from aiogram.dispatcher.filters import Text
from aiogram.dispatcher.filters.state import State, StatesGroup
from aiogram.types import InlineKeyboardMarkup, InlineKeyboardButton, ReplyKeyboardMarkup, KeyboardButton
from aiogram.utils import executor
import google.generativeai as genai
import requests
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO
import re
from PIL import Image
import tempfile
import datetime  # Qo'shilgan: kunlik limit uchun
from dotenv import load_dotenv
import os
import logging

# 1️⃣ .env faylni yuklaymiz
load_dotenv()  # Bu faylni o‘qiydi va OS environment ga qo‘shadi

# Configure logging
logging.basicConfig(filename='bot_log.txt', level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# Environment variables (Set these in PythonAnywhere environment variables)
BOT_TOKEN = os.getenv('BOT_TOKEN')        # <-- kalit nomini yozamiz
GEMINI_API_KEY = os.getenv('GEMINI_API_KEY')
PEXELS_API_KEY = os.getenv('PEXELS_API_KEY')
ADMIN_ID = int(os.getenv('ADMIN_ID', '0'))

# Fallback image URL (a generic placeholder image)
FALLBACK_IMAGE_URL = 'https://via.placeholder.com/800x600?text=No+Image+Found'

# Design styles mapping (simple color schemes for demonstration)
DESIGN_STYLES = {
    'Minimal': {'bg_color': RGBColor(255, 255, 255), 'text_color': RGBColor(0, 0, 0), 'accent_color': RGBColor(0, 0, 255)},
    'Business': {'bg_color': RGBColor(0, 51, 102), 'text_color': RGBColor(255, 255, 255), 'accent_color': RGBColor(255, 204, 0)},
    'Academic': {'bg_color': RGBColor(255, 255, 255), 'text_color': RGBColor(0, 0, 0), 'accent_color': RGBColor(128, 0, 0)},
    'Dark': {'bg_color': RGBColor(34, 34, 34), 'text_color': RGBColor(255, 255, 255), 'accent_color': RGBColor(0, 255, 0)},
    'Gradient': {'bg_color': RGBColor(0, 128, 255), 'text_color': RGBColor(255, 255, 255), 'accent_color': RGBColor(255, 0, 128)},  # Simplified, no real gradient
    'Startup': {'bg_color': RGBColor(0, 255, 204), 'text_color': RGBColor(0, 0, 0), 'accent_color': RGBColor(255, 102, 0)},
    'Elegant': {'bg_color': RGBColor(230, 230, 250), 'text_color': RGBColor(75, 0, 130), 'accent_color': RGBColor(218, 165, 32)},
    'Tech': {'bg_color': RGBColor(0, 0, 0), 'text_color': RGBColor(0, 255, 255), 'accent_color': RGBColor(255, 0, 0)},
    'Colorful': {'bg_color': RGBColor(255, 255, 0), 'text_color': RGBColor(0, 0, 0), 'accent_color': RGBColor(255, 0, 255)},
    'Modern': {'bg_color': RGBColor(245, 245, 245), 'text_color': RGBColor(33, 33, 33), 'accent_color': RGBColor(0, 188, 212)},
}

# Language translations
LANGUAGES = {
    'uz': {
        'start': "Xush kelibsiz! Ismingiz nima?",
        'presentation_button': "Presentatsiya tayyorlash",
        'choose_language': "Tilni tanlang:",
        'enter_topic': "Mavzu nomini kiriting:",
        'choose_style': "Dizayn stilini tanlang:",
        'choose_slides': "Slayd sonini tanlang:",
        'preparing': "Presentatsiya tayyorlanmoqda...",
        'error': "Xatolik yuz berdi. Iltimos, qaytadan urinib ko'ring.",
        'back': "Orqaga",
        'cancel': "Bekor qilish",
        'limit_exceeded': "Kunlik limit (50 ta) oshib ketdi. Ertaga urinib ko'ring.",
    },
    'ru': {
        'start': "Добро пожаловать! Как ваше имя?",
        'presentation_button': "Подготовить презентацию",
        'choose_language': "Выберите язык:",
        'enter_topic': "Введите название темы:",
        'choose_style': "Выберите стиль дизайна:",
        'choose_slides': "Выберите количество слайдов:",
        'preparing': "Подготавливается презентация...",
        'error': "Произошла ошибка. Пожалуйста, попробуйте снова.",
        'back': "Назад",
        'cancel': "Отмена",
        'limit_exceeded': "Дневной лимит (50) превышен. Попробуйте завтра.",
    },
    'en': {
        'start': "Welcome! What's your name?",
        'presentation_button': "Prepare Presentation",
        'choose_language': "Choose language:",
        'enter_topic': "Enter topic name:",
        'choose_style': "Choose design style:",
        'choose_slides': "Choose number of slides:",
        'preparing': "Preparing presentation...",
        'error': "An error occurred. Please try again.",
        'back': "Back",
        'cancel': "Cancel",
        'limit_exceeded': "Daily limit (50) exceeded. Try again tomorrow.",
    }
}

# FSM States
class PresentationStates(StatesGroup):
    name = State()
    start_presentation = State()
    language = State()
    topic = State()
    style = State()
    slides = State()
    generating = State()

# Bot and Dispatcher
bot = Bot(token=BOT_TOKEN)
dp = Dispatcher(bot)
dp.middleware.setup(LoggingMiddleware())

# Admin va statistika uchun global variables
users_set = set()  # Unikal foydalanuvchilar
daily_ppt_count = 0  # Kunlik PPT soni
last_date = datetime.date.today()  # Oxirgi sana
DAILY_LIMIT = 50  # Kunlik limit

# Helper functions
def check_daily_limit():
    global daily_ppt_count, last_date
    current_date = datetime.date.today()
    if current_date != last_date:
        daily_ppt_count = 0
        last_date = current_date
    return daily_ppt_count < DAILY_LIMIT

async def get_keyboard(options, row_width=2, lang='en', add_back_cancel=False):
    keyboard = InlineKeyboardMarkup(row_width=row_width)
    for opt in options:
        keyboard.insert(InlineKeyboardButton(opt, callback_data=opt))
    if add_back_cancel:
        keyboard.add(InlineKeyboardButton(LANGUAGES[lang]['back'], callback_data='back'))
        keyboard.add(InlineKeyboardButton(LANGUAGES[lang]['cancel'], callback_data='cancel'))
    return keyboard

async def generate_slide_contents(topic, num_slides, lang='en'):
    try:
        genai.configure(api_key=GEMINI_API_KEY)
        model = genai.GenerativeModel('gemini-pro')
        prompt = f"Generate {num_slides} slides for a presentation on '{topic}' in {lang.capitalize()}. Each slide should have a short title and 3-5 bullet points. Format as: Slide X: Title\n- Bullet1\n- Bullet2\n..."
        response = model.generate_content(prompt)
        slides = re.findall(r'Slide \d+: (.*?)\n((?:- .*?\n?)+)', response.text, re.DOTALL)

        contents = []
        for title, bullets in slides:
            keywords = re.findall(r'\b\w+\b', title + ' ' + bullets)[:5]  # Simple keyword extraction
            contents.append({'title': title.strip(), 'bullets': [b.strip('- ') for b in bullets.split('\n') if b.strip()], 'keywords': ' '.join(keywords)})

        logging.info(f"Generated {len(contents)} slides for topic: {topic}")
        return contents
    except Exception as e:
        logging.error(f"Gemini API error: {e}")
        return None

async def get_image_url(keywords):
    try:
        url = f"https://api.pexels.com/v1/search?query={keywords}&per_page=1"
        headers = {"Authorization": PEXELS_API_KEY}
        response = requests.get(url, headers=headers)
        if response.status_code == 200:
            data = response.json()
            if data['photos']:
                return data['photos'][0]['src']['large']
        return FALLBACK_IMAGE_URL
    except Exception as e:
        logging.error(f"Pexels API error: {e}")
        return FALLBACK_IMAGE_URL

async def generate_pptx(slides_contents, style, topic, lang):
    try:
        prs = Presentation()
        design = DESIGN_STYLES.get(style, DESIGN_STYLES['Minimal'])

        # Title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = topic
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title.text_frame.paragraphs[0].font.color.rgb = design['text_color']
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = design['bg_color']

        for content in slides_contents:
            slide_layout = prs.slide_layouts[1]  # Title and content
            slide = prs.slides.add_slide(slide_layout)
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = design['bg_color']

            # Title
            title_shape = slide.shapes.title
            title_shape.text = content['title']
            title_shape.text_frame.paragraphs[0].font.color.rgb = design['accent_color']

            # Content (bullets)
            tf = slide.shapes.placeholders[1].text_frame
            for bullet in content['bullets']:
                p = tf.add_paragraph()
                p.text = bullet
                p.alignment = PP_ALIGN.LEFT
                p.font.color.rgb = design['text_color']

            # Image
            img_url = await get_image_url(content['keywords'])
            img_response = requests.get(img_url)
            img_stream = BytesIO(img_response.content)
            img = Image.open(img_stream)
            with tempfile.NamedTemporaryFile(delete=False, suffix='.jpg') as tmp:
                img.save(tmp.name)
                left = Inches(5)
                top = Inches(2)
                width = Inches(4)
                slide.shapes.add_picture(tmp.name, left, top, width=width)
            os.unlink(tmp.name)

        pptx_stream = BytesIO()
        prs.save(pptx_stream)
        pptx_stream.seek(0)
        logging.info("PPTX generated successfully")
        return pptx_stream
    except Exception as e:
        logging.error(f"PPTX generation error: {e}")
        return None

# Handlers
@dp.message_handler(commands=['start'])
async def start(message: types.Message, state: FSMContext):
    users_set.add(message.from_user.id)  # Foydalanuvchini qo'shish
    logging.info(f"User {message.from_user.id} started bot")
    await state.set_state(PresentationStates.name.state)
    await message.reply(LANGUAGES['en']['start'])  # Default English

@dp.message_handler(state=PresentationStates.name)
async def get_name(message: types.Message, state: FSMContext):
    await state.update_data(name=message.text)
    keyboard = ReplyKeyboardMarkup(resize_keyboard=True)
    keyboard.add(KeyboardButton(LANGUAGES['en']['presentation_button']))
    await state.set_state(PresentationStates.start_presentation.state)
    await message.reply(f"Hello, {message.text}! Press the button to start.", reply_markup=keyboard)

@dp.message_handler(Text(equals=LANGUAGES['en']['presentation_button']), state=PresentationStates.start_presentation)
async def start_presentation(message: types.Message, state: FSMContext):
    lang_options = ['🇺🇿 O‘zbek', '🇷🇺 Rus', '🇬🇧 English']
    keyboard = await get_keyboard(lang_options, row_width=3)
    await state.set_state(PresentationStates.language.state)
    await message.reply(LANGUAGES['en']['choose_language'], reply_markup=keyboard)

@dp.callback_query_handler(state=PresentationStates.language)
async def choose_language(callback: types.CallbackQuery, state: FSMContext):
    lang_map = {'🇺🇿 O‘zbek': 'uz', '🇷🇺 Rus': 'ru', '🇬🇧 English': 'en'}
    lang = lang_map.get(callback.data, 'en')
    await state.update_data(lang=lang)
    await callback.message.edit_text(LANGUAGES[lang]['enter_topic'])
    await state.set_state(PresentationStates.topic.state)

@dp.message_handler(state=PresentationStates.topic)
async def get_topic(message: types.Message, state: FSMContext):
    await state.update_data(topic=message.text)
    data = await state.get_data()
    lang = data.get('lang', 'en')
    styles = list(DESIGN_STYLES.keys())
    keyboard = await get_keyboard(styles, row_width=2, lang=lang, add_back_cancel=True)
    await message.reply(LANGUAGES[lang]['choose_style'], reply_markup=keyboard)
    await state.set_state(PresentationStates.style.state)

@dp.callback_query_handler(state=PresentationStates.style)
async def choose_style(callback: types.CallbackQuery, state: FSMContext):
    data = await state.get_data()
    lang = data.get('lang', 'en')
    if callback.data == 'back':
        await start_presentation(callback.message, state)  # Go back
        return
    elif callback.data == 'cancel':
        await state.finish()
        await callback.message.edit_text("Cancelled.")
        return
    await state.update_data(style=callback.data)
    slide_options = ['5', '10', '15']
    keyboard = await get_keyboard(slide_options, row_width=3, lang=lang, add_back_cancel=True)
    await callback.message.edit_text(LANGUAGES[lang]['choose_slides'], reply_markup=keyboard)
    await state.set_state(PresentationStates.slides.state)

@dp.callback_query_handler(state=PresentationStates.slides)
async def choose_slides(callback: types.CallbackQuery, state: FSMContext):
    data = await state.get_data()
    lang = data.get('lang', 'en')
    if callback.data == 'back':
        await get_topic(callback.message, state)  # Approximate back
        return
    elif callback.data == 'cancel':
        await state.finish()
        await callback.message.edit_text("Cancelled.")
        return
    await state.update_data(slides=int(callback.data))
    await callback.message.edit_text(LANGUAGES[lang]['preparing'])
    await state.set_state(PresentationStates.generating.state)

    # Limitni tekshirish
    if not check_daily_limit():
        await callback.message.reply(LANGUAGES[lang]['limit_exceeded'])
        await state.finish()
        return

    # Async generation
    topic = data['topic']
    num_slides = data['slides']
    style = data['style']

    slides_contents = await generate_slide_contents(topic, num_slides, lang)
    if not slides_contents:
        await callback.message.reply(LANGUAGES[lang]['error'])
        await state.finish()
        return

    pptx_stream = await generate_pptx(slides_contents, style, topic, lang)
    if not pptx_stream:
        await callback.message.reply(LANGUAGES[lang]['error'])
        await state.finish()
        return

    await bot.send_document(callback.from_user.id, ('presentation.pptx', pptx_stream))
    logging.info(f"Sent PPTX to user {callback.from_user.id}")

    # Statistika yangilash
    global daily_ppt_count
    daily_ppt_count += 1
    await state.finish()

# Admin komandalari
@dp.message_handler(commands=['stats'])
async def stats(message: types.Message):
    if message.from_user.id != ADMIN_ID:
        await message.reply("Siz admin emassiz.")
        return
    check_daily_limit()  # Sana yangilash
    stats_text = f"Jami foydalanuvchilar: {len(users_set)}\nKunlik PPTlar: {daily_ppt_count}/{DAILY_LIMIT}\nOxirgi sana: {last_date}"
    await message.reply(stats_text)
    logging.info(f"Admin {ADMIN_ID} requested stats")

@dp.message_handler(commands=['logs'])
async def send_logs(message: types.Message):
    if message.from_user.id != ADMIN_ID:
        await message.reply("Siz admin emassiz.")
        return
    try:
        with open('bot_log.txt', 'rb') as log_file:
            await bot.send_document(message.from_user.id, log_file)
        logging.info(f"Admin {ADMIN_ID} requested logs")
    except Exception as e:
        await message.reply("Log fayl topilmadi.")
        logging.error(f"Logs send error: {e}")

# Error handler
@dp.errors_handler()
async def error_handler(update, error):
    logging.error(f"Error: {error}")
    return True

if __name__ == '__main__':
    executor.start_polling(dp, skip_updates=True)